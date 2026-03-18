import json
import time
import os
import requests
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import streamlit as st

# ==================== 配置部分 ====================
# 请替换为您的DeepSeek API密钥
DEEPSEEK_API_KEY = "your_deepseek_api_key_here"
DEEPSEEK_BASE_URL = "https://api.deepseek.com"

# ==================== 第一步：调用DeepSeek API生成结构化内容 ====================
def generate_ppt_content_with_deepseek(topic, slide_count=5, style="商业汇报"):
    """
    调用DeepSeek API生成PPT内容，要求返回JSON格式的结构化数据
    """
    try:
        client = OpenAI(
            api_key=DEEPSEEK_API_KEY,
            base_url=DEEPSEEK_BASE_URL
        )
        
        prompt = f"""
        你是一位专业的PPT设计专家和内容策划师。
        
        请为以下主题设计一份PPT：
        主题：{topic}
        幻灯片数量：{slide_count}页
        风格：{style}
        
        要求：
        1. 以JSON格式输出，结构如下：
        {{
            "slides": [
                {{
                    "title": "幻灯片标题",
                    "content": ["要点1", "要点2", "要点3"],
                    "slide_type": "title/content/title_only",  # 幻灯片类型建议
                    "chart_type": "bar/line/pie/none",  # 图表类型建议
                    "chart_data": {{"categories": [], "series": {{"name": []}}}}  # 示例图表数据
                }}
            ]
        }}
        
        2. 内容逻辑清晰，符合{style}场景需求
        3. 每页幻灯片包含3-5个关键要点
        4. 为至少2页幻灯片建议合适的图表类型和示例数据
        5. 内容专业、准确、有深度
        """
        
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一位专业的PPT内容设计师，擅长制作结构清晰、内容专业的演示文稿。"},
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"},
            temperature=0.7,
            max_tokens=2000
        )
        
        content = response.choices[0].message.content
        ppt_data = json.loads(content)
        return ppt_data
        
    except Exception as e:
        st.error(f"调用DeepSeek API时出错: {str(e)}")
        return None

# ==================== 第二步：处理用户上传的模板 ====================
def save_uploaded_template(uploaded_file):
    """
    保存用户上传的PPT模板文件[6](@ref)[7](@ref)
    """
    try:
        if uploaded_file is not None:
            # 创建模板保存目录
            template_dir = "user_templates"
            if not os.path.exists(template_dir):
                os.makedirs(template_dir)
            
            # 生成唯一的文件名
            timestamp = time.strftime("%Y%m%d_%H%M%S")
            filename = f"user_template_{timestamp}.pptx"
            save_path = os.path.join(template_dir, filename)
            
            # 保存文件
            with open(save_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            st.success(f"模板已保存: {filename}")
            return save_path
    except Exception as e:
        st.error(f"保存模板时出错: {str(e)}")
    return None

def get_default_template():
    """
    获取系统默认模板
    """
    # 创建一个简单的默认模板
    prs = Presentation()
    
    # 添加几种常用布局
    # 标题幻灯片布局
    title_layout = prs.slide_layouts
    
    # 标题+内容布局
    content_layout = prs.slide_layouts
    
    # 仅标题布局
    title_only_layout = prs.slide_layouts
    
    # 保存默认模板
    default_path = "default_template.pptx"
    prs.save(default_path)
    
    return default_path

# ==================== 第三步：使用python-pptx生成PPT文件 ====================
def create_ppt_from_data(ppt_data, template_path=None):
    """
    根据结构化数据和模板创建PPT文件[1](@ref)[2](@ref)[3](@ref)
    """
    try:
        # 加载模板或创建新演示文稿
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            st.info("使用自定义模板创建PPT")
        else:
            prs = Presentation()
            st.info("使用默认布局创建PPT")
        
        # 定义布局映射
        layout_mapping = {
            "title": 0,        # 标题幻灯片
            "content": 1,      # 标题+内容
            "title_only": 5,   # 仅标题
            "blank": 6         # 空白
        }
        
        # 遍历每页幻灯片数据
        for i, slide_info in enumerate(ppt_data.get("slides", [])):
            # 获取幻灯片类型，默认为内容布局
            slide_type = slide_info.get("slide_type", "content")
            layout_idx = layout_mapping.get(slide_type, 1)
            
            # 确保布局索引在有效范围内
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 1  # 回退到标题+内容布局
            
            # 添加幻灯片
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题（如果有标题占位符）
            if slide.shapes.title:
                slide.shapes.title.text = slide_info.get("title", f"幻灯片 {i+1}")
            
            # 添加内容要点（如果有内容占位符）
            content = slide_info.get("content", [])
            if content and len(slide.placeholders) > 1:
                # 尝试使用第二个占位符作为内容区域
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders
                    if hasattr(content_placeholder, 'text_frame'):
                        tf = content_placeholder.text_frame
                        tf.clear()
                        
                        # 添加内容要点
                        for point in content:
                            p = tf.add_paragraph()
                            p.text = point
                            p.level = 0
                            p.font.size = Pt(18)
                            p.font.color.rgb = RGBColor(0, 0, 0)
            
            # 添加图表（如果建议了图表类型）
            chart_type = slide_info.get("chart_type", "none")
            chart_data = slide_info.get("chart_data")
            
            if chart_type != "none" and chart_data and i > 0:
                add_chart_to_slide(slide, chart_type, chart_data, i)
        
        # 保存PPT文件
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        filename = f"AI_Generated_PPT_{timestamp}.pptx"
        prs.save(filename)
        
        return filename
        
    except Exception as e:
        st.error(f"创建PPT文件时出错: {str(e)}")
        return None

def add_chart_to_slide(slide, chart_type, chart_data, slide_index):
    """
    在幻灯片中添加图表[3](@ref)
    """
    try:
        # 设置图表位置和大小
        left = Inches(1)
        top = Inches(2) if slide.shapes.title else Inches(1)
        width = Inches(6)
        height = Inches(4)
        
        # 准备图表数据
        chart_data_obj = CategoryChartData()
        
        # 添加类别（X轴）
        if "categories" in chart_data:
            chart_data_obj.categories = chart_data["categories"]
        
        # 添加数据系列
        if "series" in chart_data:
            for series_name, series_data in chart_data["series"].items():
                chart_data_obj.add_series(series_name, series_data)
        
        # 根据图表类型添加图表
        if chart_type.lower() == "bar":
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left, top, width, height,
                chart_data_obj
            ).chart
        elif chart_type.lower() == "line":
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE,
                left, top, width, height,
                chart_data_obj
            ).chart
        elif chart_type.lower() == "pie":
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                left, top, width, height,
                chart_data_obj
            ).chart
            chart.has_legend = True
            chart.plots[0].has_data_labels = True
        
        st.success(f"在第{slide_index+1}页成功添加{chart_type}图表")
        
    except Exception as e:
        st.warning(f"添加图表时出错: {str(e)}，将继续生成PPT")

# ==================== 第四步：Streamlit Web界面 ====================
def main():
    """
    主函数：Streamlit Web界面[6](@ref)[7](@ref)[8](@ref)
    """
    st.set_page_config(
        page_title="AI PPT生成器（支持自定义模板）",
        page_icon="📊",
        layout="wide"
    )
    
    # 页面标题和描述
    st.title("🤖 AI PPT智能生成器（支持自定义模板）")
    st.markdown("""
    使用LYThink AI自动生成专业演示文稿。您可以上传自己的PPT模板，系统将基于您的模板生成PPT！
    **核心功能：**
    - 支持上传自定义PPT模板
    - AI智能生成内容
    - 自动添加图表和数据可视化
    - 多种幻灯片布局选择
    """)
    
    # 侧边栏配置
    with st.sidebar:
        st.header("⚙️ 生成设置")
        
        # 模板上传部分
        st.subheader("📁 模板设置")
        use_custom_template = st.checkbox("使用自定义模板", value=False)
        
        template_file = None
        if use_custom_template:
            template_file = st.file_uploader(
                "上传PPT模板文件",
                type=["pptx"],
                help="请上传您的PPT模板文件（.pptx格式）"
            )
            
            if template_file:
                st.info(f"已选择模板: {template_file.name}")
                # 显示模板预览
                with st.expander("模板预览信息"):
                    st.write(f"文件大小: {template_file.size / 1024:.2f} KB")
                    st.write("模板将在生成时应用")
        
        # PPT主题输入
        st.subheader("📝 内容设置")
        topic = st.text_input(
            "输入PPT主题",
            placeholder="例如：人工智能的未来发展",
            help="请输入您想要制作的PPT主题"
        )
        
        # 幻灯片数量选择
        slide_count = st.slider(
            "幻灯片数量",
            min_value=3,
            max_value=15,
            value=8,
            help="选择生成的幻灯片页数"
        )
        
        # 风格选择
        style = st.selectbox(
            "PPT风格",
            ["商业汇报", "学术演讲", "产品介绍", "项目提案", "教育培训"],
            help="选择适合您场景的PPT风格"
        )
        
        # 高级选项
        with st.expander("高级选项"):
            include_charts = st.checkbox("自动添加图表", value=True)
            chart_count = st.slider("图表数量", 0, 5, 2) if include_charts else 0
            
            # 布局偏好
            st.subheader("布局偏好")
            prefer_title_slides = st.checkbox("优先使用标题幻灯片", value=True)
            prefer_content_slides = st.checkbox("优先使用内容幻灯片", value=True)
        
        # 生成按钮
        generate_button = st.button(
            "🚀 生成PPT",
            type="primary",
            use_container_width=True
        )
    
    # 主内容区
    col1, col2 = st.columns([2, 1])
    
    with col1:
        if topic:
            st.subheader("📋 生成预览")
            st.info(f"**主题：** {topic}")
            st.info(f"**页数：** {slide_count}页")
            st.info(f"**风格：** {style}")
            
            if use_custom_template and template_file:
                st.success("✅ 将使用自定义模板生成")
            else:
                st.info("ℹ️ 将使用系统默认模板生成")
    
    with col2:
        st.subheader("📊 系统状态")
        st.metric("模板状态", "自定义" if use_custom_template and template_file else "默认")
        st.metric("生成时间", "10-30秒")
        st.metric("支持格式", ".pptx")
    
    # 生成PPT的逻辑
    if generate_button and topic:
        with st.spinner("🤖 AI正在思考并生成内容..."):
            # 步骤1：生成内容
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            status_text.text("步骤1/4：调用DeepSeek API生成内容...")
            ppt_data = generate_ppt_content_with_deepseek(topic, slide_count, style)
            progress_bar.progress(25)
            
            if ppt_data:
                status_text.text("步骤2/4：解析内容并准备PPT结构...")
                
                # 显示生成的内容预览
                with st.expander("📄 查看生成的内容结构"):
                    st.json(ppt_data)
                
                progress_bar.progress(50)
                
                # 步骤3：处理模板
                status_text.text("步骤3/4：处理模板文件...")
                template_path = None
                
                if use_custom_template and template_file:
                    # 保存用户上传的模板
                    template_path = save_uploaded_template(template_file)
                    if not template_path:
                        st.warning("自定义模板处理失败，将使用默认模板")
                        template_path = get_default_template()
                else:
                    # 使用默认模板
                    template_path = get_default_template()
                
                progress_bar.progress(75)
                
                # 步骤4：创建PPT文件
                status_text.text("步骤4/4：使用python-pptx生成PPT文件...")
                ppt_filename = create_ppt_from_data(ppt_data, template_path)
                progress_bar.progress(100)
                
                if ppt_filename:
                    status_text.text("✅ PPT生成完成！")
                    
                    # 提供下载
                    with open(ppt_filename, "rb") as file:
                        st.download_button(
                            label="📥 下载PPT文件",
                            data=file,
                            file_name=ppt_filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                    
                    # 显示成功信息
                    st.success(f"PPT已成功生成：**{ppt_filename}**")
                    st.balloons()
                    
                    # 显示生成统计
                    col_stat1, col_stat2, col_stat3 = st.columns(3)
                    with col_stat1:
                        st.metric("总页数", len(ppt_data.get("slides", [])))
                    with col_stat2:
                        chart_count = sum(1 for slide in ppt_data.get("slides", []) 
                                        if slide.get("chart_type") != "none")
                        st.metric("图表数量", chart_count)
                    with col_stat3:
                        st.metric("模板类型", "自定义" if use_custom_template else "默认")
                else:
                    st.error("PPT文件生成失败，请检查错误信息。")
            else:
                st.error("内容生成失败，请检查API密钥或网络连接。")
    
    # 使用说明
    with st.expander("📖 使用说明与技巧"):
        st.markdown("""
        ### 模板使用指南：
        
        1. **自定义模板准备**：
           - 在PowerPoint中设计您的模板，包含标题、内容等占位符
           - 保存为.pptx格式
           - 模板中的占位符将被AI生成的内容自动填充
        
        2. **模板兼容性**：
           - 支持所有标准的.pptx格式模板
           - 建议模板包含清晰的占位符布局
           - 系统会自动识别模板中的标题和内容区域
        
        3. **内容生成优化**：
           - 主题描述越详细，生成内容越精准
           - 选择合适的风格以获得最佳匹配
           - 图表数据会根据内容自动生成
        
        ### 技术特点：
        - 基于DeepSeek大模型，内容专业准确
        - 支持自定义模板，保持品牌一致性
        - 使用python-pptx库，生成标准.pptx格式
        - 自动添加数据可视化图表
        
        ### 注意事项：
        - 确保DeepSeek API密钥有效
        - 自定义模板需包含标准占位符
        - 首次生成可能需要较长时间
        - 建议在生成后检查内容准确性
        """)
    
    # 页脚
    st.markdown("---")
    st.caption("""
    **技术实现参考：** 
    - Python自动化生成PPT教程[1](@ref)
    - python-pptx自定义模板方法[2](@ref)[4](@ref)
    - Streamlit文件上传功能[6](@ref)[7](@ref)[8](@ref)
    - 更新时间：2026年3月18日
    """)

# ==================== 运行应用 ====================
if __name__ == "__main__":
    main()
